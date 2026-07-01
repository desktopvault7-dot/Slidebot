import os
import re
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from groq import Groq
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import (
    Application, CommandHandler, MessageHandler,
    CallbackQueryHandler, ContextTypes, filters
)

BOT_TOKEN = "8653002187:AAGbU2JBrj7CSjttqCU3HA4sb-5VFtS8yPY"
GROQ_API_KEY = "AQ.Ab8RN6JgVgjd4sQZibSnC1xJHa722gcJBDNehHdeqWtBVrEpyg"
PEXELS_API_KEY = "j2ABEAWO1ZzvpwIEG5ySMOsEdjJwb2lONQJNb8YzBHhifcGPdALPvbm6"

CHANNEL_USERNAME = "@slidemakernews"
CHANNEL_LINK = "https://t.me/slidemakernews"
                                                                                                                                                    
client = Groq(api_key=GROQ_API_KEY)
def search_pexels_image(query):
    headers = {
        "Authorization": PEXELS_API_KEY
    }

    r = requests.get(
        "https://api.pexels.com/v1/search",
        headers=headers,
        params={
            "query": query,
            "per_page": 1
        },
        timeout=20
    )

    data = r.json()

    if data.get("photos"):
        return data["photos"][0]["src"]["large"]

    return None
    
def download_image(url, filename):
    r = requests.get(url, timeout=30)

    with open(filename, "wb") as f:
        f.write(r.content)

    return filename
  

async def check_subscription(user_id, context):
    try:
        member = await context.bot.get_chat_member(CHANNEL_USERNAME, user_id)
        return member.status in ["member", "administrator", "creator"]
    except Exception:
        return False


async def send_join_message(update):
    kb = [
        [InlineKeyboardButton("📢 Kanalga a'zo bo'lish", url=CHANNEL_LINK)],
        [InlineKeyboardButton("✅ Tekshirish", callback_data="check_sub")]
    ]
    await update.message.reply_text(
        "🔒Botdan foydalanish uchun kanalga a'zo bo'ling.",
        reply_markup=InlineKeyboardMarkup(kb)
    )


async def check_button(update, context):
    q = update.callback_query
    await q.answer()

    if await check_subscription(q.from_user.id, context):
        await q.message.edit_text("✅ A'zolik tasdiqlandi.\n\n/start yuboring.")
    else:
        await q.answer("❌ Hali a'zo emassiz.", show_alert=True)


async def start(update, context):
    if not await check_subscription(update.effective_user.id, context):
        await send_join_message(update)
        return

    kb = [[InlineKeyboardButton("📊 Taqdimot yasash", callback_data="create_presentation")]]
    text = """
🎓 <b>Slide Maker Bot</b>

🤖 Sun'iy intellekt yordamida bir necha daqiqada professional taqdimot tayyorlang.

📌 Imkoniyatlar:
• 🇺🇿 O'zbek
• 🇬🇧 English
• 🇷🇺 Русский
• 🎨 Premium dizayn
• 🖼️ Avtomatik rasmlar

👇 Boshlash uchun quyidagi tugmani bosing.
"""

await update.message.reply_text(
    text,
    parse_mode="HTML",
    reply_markup=InlineKeyboardMarkup(kb)
)


async def menu_buttons(update, context):
    q = update.callback_query
    await q.answer()

    if q.data == "create_presentation":
        kb = [
            [InlineKeyboardButton("🇺🇿 O'zbek tili", callback_data="lang_uz")],
            [InlineKeyboardButton("🇬🇧 English", callback_data="lang_en")],
            [InlineKeyboardButton("🇷🇺 Русский", callback_data="lang_ru")]
        ]
        await q.message.edit_text("🌐 Taqdimot tilini tanlang:", reply_markup=InlineKeyboardMarkup(kb))

    elif q.data.startswith("lang_"):
        context.user_data.clear()
        context.user_data["language"] = q.data.replace("lang_", "")
        context.user_data["step"] = "topic"
        await q.message.edit_text("📚 Taqdimot mavzusini bexato kiriting:")
    elif q.data.startswith("template_"):

        context.user_data["template"] = q.data.replace("template_", "")

        context.user_data["step"] = "student"

        await q.message.reply_text(
        "Ismingizni kiriting.\n(Diqqat: bu ism taqdimotga kiritiladi)"
        )


def add_center_text(slide, text, size):
    box = slide.shapes.add_textbox(
        Inches(0.8),
        Inches(1.8),
        Inches(11.8),
        Inches(2.5)
    )

    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def add_background(slide, image_file):
    slide.shapes.add_picture(
        image_file,
        0,
        0,
        width=Inches(13.333),
        height=Inches(7.5)
    )


def build_presentation(topic, student, teacher, slide_count, content, lang, template):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    bg_file = f"template{template}.jpg"

    content_slides = max(1, slide_count - 2)

    # titul
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bg_file)
    add_center_text(slide, topic.upper(), 88)
    if lang == "uz":
        student_label = "Tayyorlovchi"
        teacher_label = "O'qituvchi"
    elif lang == "en":
        student_label = "Prepared by"
        teacher_label = "Teacher"
    else:
        student_label = "Подготовил"
        teacher_label = "Преподаватель"

    info = slide.shapes.add_textbox(
        Inches(9.0),
        Inches(6.1),
        Inches(4),
        Inches(0.7)
    )

    tf = info.text_frame
    tf.text = f"{student_label}: {student}\n{teacher_label}: {teacher}"

    for p in tf.paragraphs:
        p.alignment = PP_ALIGN.RIGHT
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(255, 255, 255)

    slides = re.findall(
        r"===SLIDE===\s*Title:\s*(.*?)\s*Keyword:\s*(.*?)\s*Content:\s*(.*?)(?=(===SLIDE===|$))",
        content,
        re.S
    )        
    if len(slides) < content_slides:
        raise Exception(
            f"Yetarli slayd topilmadi. Kerak: {content_slides}, topildi: {len(slides)}"
    ) 
    for title, keyword, body, _ in slides[:content_slides]:
        s = prs.slides.add_slide(prs.slide_layouts[6])
        add_background(s, bg_file)
       
        title_box = s.shapes.add_textbox(
            Inches(0.5),
            Inches(0.3),
            Inches(12.3),
            Inches(0.8)
        )

        p = title_box.text_frame.paragraphs[0]
        p.text = title.strip()
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 255)

        keyword = keyword.strip()

        img_url = search_pexels_image(keyword)

        if img_url:
            try:
                img_file = f"{keyword[:20]}.jpg"

                download_image(img_url, img_file)

                s.shapes.add_picture(
                    img_file,
                    Inches(0.4),
                    Inches(1.2),
                    width=Inches(4.2),
                    height=Inches(4.5)
                )
            except Exception as e:
                print("IMAGE ERROR:", e)

        body_box = s.shapes.add_textbox(
            Inches(5.0),
            Inches(1.2),
            Inches(7.0),
            Inches(4.8)
        )

        tf = body_box.text_frame
        tf.word_wrap = True
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        p = tf.paragraphs[0]
        p.text = body.strip()
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)

    if lang == "uz":
        thanks_text = "E'TIBORINGIZ UCHUN RAHMAT!"
    elif lang == "en":
        thanks_text = "THANK YOU FOR YOUR ATTENTION!"
    else:
        thanks_text = "СПАСИБО ЗА ВНИМАНИЕ!"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, bg_file)
    add_center_text(slide, thanks_text, 72)

    return prs


async def handle_message(update, context):
    if not await check_subscription(update.effective_user.id, context):
        await send_join_message(update)
        return

    step = context.user_data.get("step")

    if step == "topic":
        context.user_data["topic"] = update.message.text
        context.user_data["step"] = "template"

        kb = [
            [
                InlineKeyboardButton("1", callback_data="template_1"),
                InlineKeyboardButton("2", callback_data="template_2"),
                InlineKeyboardButton("3", callback_data="template_3")
            ],
            [
                InlineKeyboardButton("4", callback_data="template_4"),
                InlineKeyboardButton("5", callback_data="template_5"),
                InlineKeyboardButton("6", callback_data="template_6")
            ]
        ]

        with open("preview.jpg", "rb") as photo:
            await update.message.reply_photo(
                photo=photo,
                caption="🎨 Taqdimot uchun mavzu tanlang:",
                reply_markup=InlineKeyboardMarkup(kb)
            )
        return
        
    if step == "student":
        context.user_data["student"] = update.message.text
        context.user_data["step"] = "teacher"
        await update.message.reply_text("O'qituvchi ism-familiyasini kiriting.")
        return

    if step == "teacher":
        context.user_data["teacher"] = update.message.text
        context.user_data["step"] = "count"
        await update.message.reply_text("📄 Taqdimot necha sahifalik bo'lsin?")
        return

    if step == "count":
        try:
            slide_count = int(update.message.text)
            if slide_count < 3:
                await update.message.reply_text("Kamida 3 ta slayd kiriting.")
                return
        except:
            await update.message.reply_text("Slayd sonini raqamda yuboring.")
            return

        topic = context.user_data["topic"]
        student = context.user_data["student"]
        teacher = context.user_data["teacher"]
        lang = context.user_data["language"]

        if lang == "uz":
            instruction = "Write only in Uzbek."
        elif lang == "en":
            instruction = "Write only in English."
        else:
            instruction = "Пиши только на русском языке."

        await update.message.reply_text("⏳ Taqdimot tayyorlanmoqda...")

        content_count = slide_count - 2

        prompt = f"""
{instruction}

Topic: {topic}

Generate exactly {content_count} presentation slides.

Each slide must have:
- Unique section title
- 130-150 words of detailed educational content
- Continuous paragraph form
- No bullet points
- No numbering
- No outline style

The content should be ready to place directly into a presentation.
Keyword rules:

- Keyword must describe a visible object, place, person or scene.
- Use only 2-4 English words.
- Do not use abstract words such as:
  challenges, opportunities, importance, future, development, history, process.
- The keyword must describe what should actually appear in the image.
- Prefer real-world objects and scenes.

Examples:

Bees -> honey bee

Textile industry -> textile factory

Uzbekistan education -> school classroom

Ancient history -> ancient manuscript

Computer networks -> server room

Solar system -> planets space

Agriculture -> wheat field

Tourism -> samarkand registan

For topics related to Uzbekistan, use keywords connected to Uzbekistan, Uzbek people, Samarkand, Bukhara, Tashkent or Uzbek culture when appropriate.

Format exactly like this:

===SLIDE===
Title: Introduction to Taxi Drivers
Keyword: city traffic
Content: Full paragraph here...

===SLIDE===
Title: Responsibilities of Taxi Drivers
Keyword: taxi driver
Content: Full paragraph here...

IMPORTANT:
- Keyword must be on its own line.
- Keyword must come BEFORE Content.
- Never place Keyword inside Content.
"""

        response = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[{"role": "user", "content": prompt}]
        )

        result = response.choices[0].message.content
        print(result)

        try:
            prs = build_presentation(
                topic,
                student,
                teacher,
                slide_count,
                result,
                lang,
                context.user_data["template"]
            )

            filename = f"{update.effective_chat.id}.pptx"
            prs.save(filename)

            with open(filename, "rb") as f:
                await update.message.reply_document(
                    f,
                    filename="presentation.pptx"
                )

            os.remove(filename)
            context.user_data.clear()

        except Exception as e:
            await update.message.reply_text(f"XATO: {str(e)}")
            print(f"XATO: {e}")


def main():
    app = Application.builder().token(BOT_TOKEN).build()

    app.add_handler(CommandHandler("start", start))
    app.add_handler(CallbackQueryHandler(check_button, pattern="check_sub"))
    app.add_handler(CallbackQueryHandler(menu_buttons, pattern="^(create_presentation|lang_|template_)"))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))

    app.run_polling()


if __name__ == "__main__":
    main()
